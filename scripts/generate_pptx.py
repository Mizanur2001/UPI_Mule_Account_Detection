"""
FinGuard CSIC 1.0 PPT - Exact Reference Format (v2 - reviewed & perfected)
===========================================================================
Matches the CSIC 1.0 reference PDF template exactly:
- Aptos Display / Aptos fonts
- #0e2841 navy text color, white background
- CSIC logo at bottom-right on every slide
- Slide numbers on content slides (2-12)
- All 10 diagrams embedded where relevant
- Differentiated slide titles for Part 1/Part 2 slides
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DIAGRAMS = os.path.join(BASE_DIR, "docs", "diagrams")
REF_IMGS = os.path.join(BASE_DIR, "docs", "ref_images")
OUT_PATH = os.path.join(BASE_DIR, "FINAL_PRESENTATION.pptx")

# ── COLORS (from reference PDF) ─────────────────────────────────
NAVY = RGBColor(0x0E, 0x28, 0x41)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PX2IN = 13.333 / 960.0


def px(val):
    return Inches(val * PX2IN)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_name='Aptos',
                 font_size=Pt(12), color=NAVY, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_multi_para(slide, left, top, width, height, lines):
    """lines: list of (text, font_name, font_size, color, bold, italic, align, space_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fname, fsize, color, bold, italic, align, sp) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = fname
        p.font.size = fsize
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        if sp:
            p.space_after = sp
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=Pt(14),
                color=NAVY, font_name='Aptos', spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'\u2022  {item}'
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = spacing
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_csic_logo(slide):
    logo_path = os.path.join(REF_IMGS, "ref_img_0.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, px(870), px(450), px(70), px(75))


def add_slide_number(slide, num):
    """Small slide number at bottom-center, matching ref minimal style."""
    add_text_box(slide, Inches(6.3), px(515), Inches(0.8), px(20),
                 str(num), font_name='Aptos', font_size=Pt(10),
                 color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def add_slide_title(slide, title_text):
    add_text_box(slide, px(20), px(20), Inches(12), Inches(0.8),
                 title_text, font_name='Aptos Display', font_size=Pt(36),
                 color=NAVY)
    # Thin accent line under title
    add_rect(slide, px(20), px(65), Inches(2), Pt(2), NAVY)


def add_image_safe(slide, filename, left, top, width=None, height=None):
    path = os.path.join(DIAGRAMS, filename)
    if os.path.exists(path):
        kwargs = {}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        return slide.shapes.add_picture(path, left, top, **kwargs)
    return None


def add_table(slide, left, top, width, rows_data, font_size=Pt(10),
              header_bg=NAVY, header_fg=WHITE, row_height=Inches(0.30)):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                        width, row_height * n_rows)
    tbl = tbl_shape.table
    for r, row_data in enumerate(rows_data):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Reduce cell margins for tighter fit
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for para in cell.text_frame.paragraphs:
                para.font.name = 'Aptos'
                para.font.size = font_size
                para.alignment = PP_ALIGN.CENTER
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = header_fg
                else:
                    para.font.color.rgb = NAVY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY
    return tbl_shape


# ===================================================================
# BUILD ALL 12 SLIDES
# ===================================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ==============================================================
    # SLIDE 1: TITLE
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)

    # Top navy banner
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, px(75), NAVY)
    add_text_box(sl, px(177), px(18), Inches(9), px(50),
                 'Cyber Security Innovation Challenge 1.0',
                 font_name='Aptos Display', font_size=Pt(36),
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # PROTOTYPE DEVELOPMENT centered
    add_text_box(sl, Inches(0), px(100), SLIDE_W, Inches(1.2),
                 'PROTOTYPE\nDEVELOPMENT',
                 font_name='Aptos Display', font_size=Pt(54),
                 color=NAVY, alignment=PP_ALIGN.CENTER)

    # Stage - III
    add_text_box(sl, Inches(0), px(230), SLIDE_W, Inches(0.9),
                 'Stage \u2013 III',
                 font_name='Aptos Display', font_size=Pt(50),
                 color=NAVY, alignment=PP_ALIGN.CENTER)

    # Team info - filled with actual data from user's PDF
    fields = [
        'Team Name \u2013 MandelBrot',
        'Team Lead Name \u2013 Mizanur Rahaman',
        'Team Lead University \u2013 Jadavpur University',
        'Problem Statement Domain \u2013 Mule Accounts & Collusive Fraud in UPI',
        'Solution Subtitle \u2013 FinGuard: Real-Time Multi-Signal Mule Account Detection',
    ]
    y = px(300)
    for f in fields:
        add_text_box(sl, px(127), y, Inches(10), px(35),
                     f, font_name='Aptos', font_size=Pt(24),
                     color=BLACK)
        y += px(33)

    add_csic_logo(sl)

    # ==============================================================
    # SLIDE 2: PROBLEM STATEMENT & CONTEXT
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Problem Statement & Context')

    # Left column - Crisis
    add_text_box(sl, px(25), px(80), Inches(6), px(28),
                 'The Mule Account Crisis in UPI',
                 font_name='Aptos Display', font_size=Pt(18),
                 color=NAVY, bold=True)

    add_bullets(sl, px(25), px(112), Inches(6), Inches(3.5), [
        'UPI processed 13.1 billion transactions in Oct 2024 (\u20b920.64 lakh crore)',
        'Mule accounts: bank accounts used as pass-throughs for laundering stolen money',
        'Individual transactions look normal; fraud is in the coordination and network patterns',
        'Current rule engines cannot see networks, go stale, and offer no nuanced scoring',
        'RBI mandates enhanced fraud monitoring; NPCI requires multi-dimensional detection',
    ], font_size=Pt(13))

    # Right column - Formulation
    add_text_box(sl, px(500), px(80), Inches(6), px(28),
                 'Our Formulation',
                 font_name='Aptos Display', font_size=Pt(18),
                 color=NAVY, bold=True)

    add_bullets(sl, px(500), px(112), Inches(6), Inches(3.5), [
        'Build transaction graph G = (A, E) from UPI account data',
        'Five-signal ensemble: Behavioral (25%), Graph (40%), Device (15%), Temporal (10%), ML (10%)',
        'R(a) = \u03a3 w\u2096 \u00b7 S\u2096(a) + Boost({S\u2096(a)}) -- composite risk score per account',
        'Confidence boosting: +8 to +20 when 2-4 independent signals agree',
        'Risk tiers: CRITICAL (\u226585), HIGH (70-84), MEDIUM (40-69), LOW (<40)',
        'Target: sub-50ms per-account scoring latency for real-time UPI integration',
    ], font_size=Pt(13))

    # Bottom stats bar
    add_rect(sl, px(20), px(448), Inches(12.8), Pt(1.5), NAVY)
    add_text_box(sl, px(25), px(453), Inches(12.5), px(35),
                 '13.1B txns/month   |   \u20b920.64L Cr volume   |   300+ banks on UPI   |   5 detection signals   |   <50ms scoring',
                 font_name='Aptos', font_size=Pt(14),
                 color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    add_csic_logo(sl)
    add_slide_number(sl, 2)

    # ==============================================================
    # SLIDE 3: REVIEW OF EXISTING SOLUTIONS
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Review of Existing Solutions and Research')

    gap_data = [
        ['Capability', 'Rule-Based', 'Supervised ML', 'Graph Methods', 'FinGuard (Ours)'],
        ['Temporal Pattern Detection', 'Static thresholds', 'Limited', 'None', 'Full (5 sub-signals)'],
        ['Network/Graph Analysis', 'None', 'None', 'Yes', 'Yes (3 patterns + DFS)'],
        ['Device Correlation', 'None', 'Partial', 'None', 'Full'],
        ['Unsupervised (No labels)', 'N/A', 'No', 'Partial', 'Yes (IF + Z-score)'],
        ['Real-Time Scoring', 'Fast', 'Moderate', 'Slow', 'Fast (<50ms)'],
        ['Explainability', 'Clear', 'Black-box', 'Limited', 'Full (3-5 evidence)'],
        ['Multi-Signal Ensemble', 'No', 'No', 'No', 'Yes (5-factor weighted)'],
        ['Confidence Boosting', 'No', 'No', 'No', 'Yes (multi-signal)'],
    ]
    add_table(sl, px(20), px(78), Inches(7.2), gap_data, font_size=Pt(9),
              row_height=Inches(0.33))

    add_text_box(sl, px(560), px(78), Inches(5.5), px(25),
                 'Key Research',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_bullets(sl, px(560), px(108), Inches(5.5), Inches(3.5), [
        'MuleTrack (Jambhrunkar 2025): Lightweight temporal learning for mule detection',
        'GNN Review (Cheng 2024): Graph methods beat classifiers for coordinated fraud',
        'Node2Vec (Caglayan 2022): Graph embeddings improve laundering detection',
        'Community Detection (Huang 2025): Mules found via graph communities',
        'Isolation Forest (Liu 2008): Unsupervised anomaly detection without labels',
        'Neo4j (2023): Industry case studies on graph-based fraud detection',
    ], font_size=Pt(11))

    add_text_box(sl, px(20), px(470), Inches(12.5), px(35),
                 'Gap: No existing system combines all 5 signals (behavioral + graph + device + temporal + ML) with explainability and confidence boosting. FinGuard fills this gap.',
                 font_name='Aptos', font_size=Pt(13),
                 color=NAVY, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    add_csic_logo(sl)
    add_slide_number(sl, 3)

    # ==============================================================
    # SLIDE 4: PROPOSED SOLUTION - ARCHITECTURE & APPROACH
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Proposed Solution \u2013 Architecture & Approach')

    # Architecture diagram on left
    add_image_safe(sl, '02_architecture.png', px(15), px(78), width=Inches(6.2))

    # Detection modules on right
    add_text_box(sl, px(475), px(78), Inches(6.2), px(24),
                 'Five Detection Modules',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_bullets(sl, px(475), px(106), Inches(6.2), Inches(2.2), [
        'Behavioral (25%): Velocity, flow asymmetry, amount anomalies, new account flags',
        'Graph Analytics (40%): Star, chain, circular patterns. Custom DFS/BFS. O(V\u00b7d)',
        'Device Fingerprinting (15%): Multi-account device sharing + device rotation',
        'Temporal (10%): Burst detection, night activity, velocity spikes, bot signatures',
        'ML Anomaly (10%): Custom Isolation Forest (NumPy), 17 features, Z-score ensemble',
    ], font_size=Pt(12))

    # Risk aggregation
    add_text_box(sl, px(475), px(275), Inches(6.2), px(24),
                 'Risk Aggregation Formula',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_text_box(sl, px(475), px(303), Inches(6.2), px(28),
                 'R(a) = 0.25\u00b7S_B + 0.40\u00b7S_G + 0.15\u00b7S_D + 0.10\u00b7S_T + 0.10\u00b7S_ML + Boost',
                 font_name='Aptos', font_size=Pt(12),
                 color=NAVY, italic=True)

    # Boost table
    boost = [
        ['Condition', 'Boost'],
        ['\u22654 signals above threshold', '+20'],
        ['\u22653 signals above threshold', '+15'],
        ['\u22652 signals above threshold', '+8'],
        ['Graph \u2265 30 AND Device \u2265 15', '+10'],
        ['Behav \u2265 40 AND Graph \u2265 40 AND Device \u2265 30', '+12'],
    ]
    add_table(sl, px(475), px(340), Inches(6), boost, font_size=Pt(9),
              row_height=Inches(0.26))

    add_csic_logo(sl)
    add_slide_number(sl, 4)

    # ==============================================================
    # SLIDE 5: INNOVATION AND NOVELTY
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Innovation and Novelty Elements')

    innovations = [
        ('1. Five-Signal Ensemble',
         'Combines behavioral + graph + device + temporal + ML into a single weighted score. Evading all 5 signals simultaneously is near-impossible. Existing systems use 1-2 signals at most.'),
        ('2. Multi-Signal Confidence Boosting',
         'When 2-4 independent signals agree, boosting adds +8 to +20 points. Corroboration from separate detection methods provides strong evidential weight.'),
        ('3. Zero-Label ML Detection',
         'Custom Isolation Forest in pure NumPy (~200 lines). No labelled training data needed. Deploys on any new platform and starts flagging outliers from day one.'),
        ('4. Efficient Graph Algorithms',
         'Custom DFS cycle detector with depth cap of 6 replaces exponential nx.simple_cycles(). BFS chain detection. O(V\u00b7d) time complexity.'),
        ('5. Explainability by Design',
         'Every module generates 3-5 specific evidence items as core logic, not a post-hoc add-on. Investigators see exactly why an account was flagged.'),
        ('6. Production-Grade Security',
         'API-key auth, rate limiting (120/min), restricted CORS, JSON audit logs with request IDs, non-root Docker container. Banking-system ready.'),
    ]

    # Signal radar diagram on the right
    add_image_safe(sl, '10_signal_radar.png', px(700), px(80), width=Inches(3.3))

    y = px(82)
    for title, desc in innovations:
        add_text_box(sl, px(30), y, Inches(9), px(16),
                     title, font_name='Aptos', font_size=Pt(14),
                     color=NAVY, bold=True)
        add_text_box(sl, px(48), y + px(19), Inches(8.8), px(42),
                     desc, font_name='Aptos', font_size=Pt(11),
                     color=NAVY)
        y += px(68)

    add_csic_logo(sl)
    add_slide_number(sl, 5)

    # ==============================================================
    # SLIDE 6: USP - PROPOSED SOLUTION VS EXISTING
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)

    add_text_box(sl, px(20), px(12), Inches(12), px(50),
                 'Unique Selling Proposition (USP): Proposed Solution\nvs. Existing Solutions (Relevance to Industry)',
                 font_name='Aptos Display', font_size=Pt(30),
                 color=NAVY)

    # Thin accent line
    add_rect(sl, px(20), px(68), Inches(2), Pt(2), NAVY)

    # Left - Key differentiators
    add_text_box(sl, px(25), px(95), Inches(6), px(24),
                 'Key Differentiators',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_bullets(sl, px(25), px(122), Inches(6), Inches(2.5), [
        '5 signals vs. 1-2: catches accounts that single-model systems miss entirely',
        'Graph-first (40% weight): reveals star, chain, and loop patterns of organized rings',
        'No training data required: Isolation Forest runs unsupervised from day one',
        'Full explainability: component breakdowns + evidence + confidence for every score',
    ], font_size=Pt(12))

    # Right - Competitive table
    comp_data = [
        ['Feature', 'Static Rules', 'ML Model', 'FinGuard'],
        ['Detection Signals', '1 (rules)', '1 (features)', '5 (ensemble)'],
        ['Graph Awareness', 'None', 'None', 'Full'],
        ['Device Correlation', 'None', 'Partial', 'Full'],
        ['Labels Required', 'No', 'Yes', 'No'],
        ['Explainability', 'High', 'Low', 'High'],
        ['Real-Time', 'Yes', 'Moderate', '<50ms'],
        ['Confidence Levels', 'No', 'No', 'Yes'],
    ]
    add_table(sl, px(500), px(95), Inches(6), comp_data, font_size=Pt(10),
              row_height=Inches(0.28))

    # Bottom - Business model
    add_rect(sl, px(20), px(318), Inches(12.8), Pt(1), NAVY)

    add_text_box(sl, px(25), px(325), Inches(12.5), px(24),
                 'Business Model & Market Viability',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_bullets(sl, px(25), px(352), Inches(6), Inches(2), [
        'Target: 300+ banks and 50+ PSPs on UPI, most using static rule engines',
        'Market: fraud losses in thousands of crores annually',
        'Each investigation costs \u20b915K-25K in analyst time',
    ], font_size=Pt(11))

    add_bullets(sl, px(500), px(352), Inches(6), Inches(2), [
        'Starter: \u20b950K/mo | Enterprise: \u20b93-5 lakh/mo with SLA',
        'SaaS API for small banks; on-premise Docker/K8s for NPCI',
        'Infrastructure licensing for large-scale deployments',
    ], font_size=Pt(11))

    add_csic_logo(sl)
    add_slide_number(sl, 6)

    # ==============================================================
    # SLIDE 7: PROTOTYPE - TECH STACK & TEST SCENARIOS
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Prototype: Technology Stack & Test Scenarios')

    # Tech stack table
    add_text_box(sl, px(25), px(78), Inches(5), px(20),
                 'Technology Stack',
                 font_name='Aptos Display', font_size=Pt(15),
                 color=NAVY, bold=True)

    tech = [
        ['Component', 'Technology', 'Version'],
        ['Backend API', 'FastAPI + Uvicorn', '2.1.0'],
        ['Frontend', 'React + Vite', '18.x / 5.x'],
        ['Language', 'Python', '3.11'],
        ['Graph Engine', 'NetworkX', '3.2.1'],
        ['ML Engine', 'Custom Isolation Forest', 'NumPy 1.26'],
        ['Containers', 'Docker + Compose', 'Multi-stage'],
    ]
    add_table(sl, px(25), px(100), Inches(5.5), tech, font_size=Pt(10))

    # Test scenarios
    add_text_box(sl, px(490), px(78), Inches(6.2), px(20),
                 'Test Scenarios (6 Mule Patterns)',
                 font_name='Aptos Display', font_size=Pt(15),
                 color=NAVY, bold=True)

    scenarios = [
        ['Scenario', 'Pattern', 'Expected Risk'],
        ['Star Aggregator', '5\u21921 mule\u21921 dist.\u21923 sinks', 'CRITICAL/HIGH'],
        ['Circular Network', '4-node loop + shared device', 'CRITICAL'],
        ['Chain Laundering', '5-node sequential chain', 'HIGH'],
        ['Device Ring', '3 accounts, 1 shared device', 'HIGH/MEDIUM'],
        ['Rapid Onboarding', '1-day acct, 13 txns in 30min', 'CRITICAL'],
        ['Night Smurfing', '12+ txns between 1-4 AM', 'HIGH'],
    ]
    add_table(sl, px(490), px(100), Inches(6.5), scenarios, font_size=Pt(10))

    # System flow + graph patterns diagrams
    add_image_safe(sl, '01_system_flow.png', px(15), px(305), width=Inches(6))
    add_image_safe(sl, '03_graph_patterns.png', px(480), px(305), width=Inches(6.5))

    add_csic_logo(sl)
    add_slide_number(sl, 7)

    # ==============================================================
    # SLIDE 8: PROTOTYPE - DETECTION RESULTS & SECURITY
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Prototype: Detection Results & Security Validation')

    # Detection results table (left)
    add_text_box(sl, px(25), px(78), Inches(6.5), px(20),
                 'Detection Results: 100% Mule Detection, 0% False Positives',
                 font_name='Aptos Display', font_size=Pt(14),
                 color=NAVY, bold=True)

    results = [
        ['Scenario', 'Score', 'Level', 'Primary Evidence'],
        ['Star Aggregator', '92', 'CRITICAL', 'Star pattern (5\u21921 out)'],
        ['Circular Network', '88', 'CRITICAL', 'Cycle + shared device'],
        ['Chain Laundering', '76', 'HIGH', 'Deep chain (4+ hops)'],
        ['Device Ring', '72', 'HIGH', 'Shared device (3 accts)'],
        ['Rapid Onboarding', '95', 'CRITICAL', 'Burst + 1-day account'],
        ['Night Smurfing', '78', 'HIGH', 'Night activity (85%)'],
    ]
    add_table(sl, px(25), px(100), Inches(6.2), results, font_size=Pt(10))

    # Security testing table (right)
    add_text_box(sl, px(490), px(78), Inches(6.2), px(20),
                 'Security Testing: 42/42 Tests Passed',
                 font_name='Aptos Display', font_size=Pt(14),
                 color=NAVY, bold=True)

    sec = [
        ['Category', 'Tests', 'Passed', 'Result'],
        ['API Key Authentication', '6', '6', '100%'],
        ['Injection Attacks', '10', '10', '100%'],
        ['Rate Limiting', '3', '3', '100%'],
        ['CORS Policy', '3', '3', '100%'],
        ['Input Validation', '7', '7', '100%'],
        ['HTTP Method Restriction', '4', '4', '100%'],
        ['Security Headers', '5', '5', '100%'],
        ['Audit Logging', '4', '4', '100%'],
        ['Total', '42', '42', '100%'],
    ]
    add_table(sl, px(490), px(100), Inches(6.2), sec, font_size=Pt(9),
              row_height=Inches(0.26))

    # Performance metrics (bottom-left)
    add_text_box(sl, px(25), px(318), Inches(4), px(20),
                 'Performance Metrics',
                 font_name='Aptos Display', font_size=Pt(14),
                 color=NAVY, bold=True)

    perf = [
        ['Metric', 'Value'],
        ['Single account scoring', '<50ms'],
        ['Batch (30 accounts)', '<500ms'],
        ['API startup', '<3s'],
        ['Memory footprint', '<150MB'],
    ]
    add_table(sl, px(25), px(340), Inches(3.8), perf, font_size=Pt(10))

    # Risk distribution diagram (bottom-center)
    add_image_safe(sl, '06_risk_distribution.png', px(310), px(318), width=Inches(3.5))

    # Security results diagram (bottom-right)
    add_image_safe(sl, '09_security_results.png', px(590), px(318), width=Inches(4))

    add_csic_logo(sl)
    add_slide_number(sl, 8)

    # ==============================================================
    # SLIDE 9: LIMITATIONS AND CONSTRAINTS
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Limitations and Constraints')

    # Left - Limitations
    add_text_box(sl, px(25), px(82), Inches(6), px(22),
                 'Current Limitations',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_bullets(sl, px(25), px(108), Inches(6), Inches(3.8), [
        'Synthetic Data: All testing on generated data; real-world patterns are messier and noisier',
        'Static Graph: Built once at startup from CSV; production needs incremental streaming updates',
        'In-Memory Only: Everything in RAM; will not scale to UPI billions of transactions natively',
        'No GNNs: Hand-crafted features; Graph Neural Networks could learn subtler patterns',
        'Fixed Weights: Ensemble weights set manually (25/40/15/10/10), not learned from data',
        'Batch ML: Isolation Forest trains once at startup; needs online/incremental retraining',
        'Single-Hop Devices: Direct sharing only, no transitive multi-hop device chains',
    ], font_size=Pt(12))

    # Right - Challenges + device flow diagram
    add_text_box(sl, px(500), px(82), Inches(6), px(22),
                 'Challenges Encountered',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    add_bullets(sl, px(500), px(108), Inches(6), Inches(2.5), [
        'Cycle Detection: nx.simple_cycles() hung on dense graphs; custom DFS with depth cap solved it',
        'Ensemble Calibration: Finding right signal weights required iterative experiments',
        'Explainability vs. Privacy: Evidence items raise questions about information surfacing',
        'Cross-Platform: Windows dev \u2192 Linux Docker deployment caused path and encoding issues',
    ], font_size=Pt(12))

    # Device flow diagram on bottom-right
    add_image_safe(sl, '04_device_flow.png', px(520), px(290), width=Inches(5.5))

    # Bottom note
    add_text_box(sl, px(25), px(450), Inches(12.5), px(35),
                 'Every limitation has a concrete production solution mapped in our 16-week MVP roadmap. These are engineering challenges with known solutions, not fundamental design flaws.',
                 font_name='Aptos', font_size=Pt(13),
                 color=NAVY, italic=True, alignment=PP_ALIGN.CENTER)

    add_csic_logo(sl)
    add_slide_number(sl, 9)

    # ==============================================================
    # SLIDE 10: ROADMAP TO MVP
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Roadmap and Strategy Towards MVP')

    phases = [
        ('Phase 1: Infrastructure\n(Weeks 1-4)', [
            'Apache Kafka for real-time event ingestion',
            'PostgreSQL for persistent transaction storage',
            'Neo4j graph DB with incremental updates',
            'Redis for sub-ms hot data cache',
        ]),
        ('Phase 2: Advanced Detection\n(Weeks 5-8)', [
            'GNN-based node classification',
            'Online/incremental anomaly detection',
            'Bidirectional graph analysis',
            'Multi-hop transitive device chains',
        ]),
        ('Phase 3: Scale & Integration\n(Weeks 9-12)', [
            'Kubernetes with horizontal autoscaling',
            'UPI switch plugin for inline scoring',
            'Case management workflow for investigators',
            'Feedback loop for model retraining',
        ]),
        ('Phase 4: Production Hardening\n(Weeks 13-16)', [
            'A/B testing framework (shadow mode)',
            'Automated SAR generation (RBI compliance)',
            'SOC 2 Type II certification',
            'Sub-10ms scoring at 10,000 TPS',
        ]),
    ]

    col_w = Inches(3.1)
    for i, (title, items) in enumerate(phases):
        x = px(20) + i * (col_w + Inches(0.12))
        y = px(85)

        add_text_box(sl, x, y, col_w, px(40),
                     title, font_name='Aptos', font_size=Pt(13),
                     color=NAVY, bold=True)

        add_bullets(sl, x, y + px(50), col_w, Inches(2),
                    items, font_size=Pt(11))

    # Roadmap timeline diagram
    add_image_safe(sl, '08_roadmap.png', px(20), px(295), width=Inches(12.5))

    add_csic_logo(sl)
    add_slide_number(sl, 10)

    # ==============================================================
    # SLIDE 11: TEAM COMPOSITION
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'Team Composition and Individual Contributions')

    members = [
        ('Team Leader', 'Mizanur Rahaman', 'Jadavpur University', '[DOB]',
         'System architecture, risk engine, ensemble calibration, final report'),
        ('Member 1', '[Name]', '[University]', '[DOB]',
         'Graph analysis, behavioral analysis, DFS cycle detection, BFS chains'),
        ('Member 2', '[Name]', '[University]', '[DOB]',
         'Custom Isolation Forest (NumPy), 17-feature pipeline, Z-score ensemble'),
        ('Member 3', '[Name]', '[University]', '[DOB]',
         'React dashboard (8 tabs), network graph vis, API integration, UX'),
        ('Member 4', '[Name]', '[University]', '[DOB]',
         'Docker containerization, security middleware, data generation, temporal analysis'),
    ]

    # Grid: row 1 = 3 across, row 2 = 2 centered (matching ref)
    positions_row1 = [(29, 91), (351, 91), (673, 91)]
    positions_row2 = [(181, 286), (526, 286)]

    for idx, (role, name, uni, dob, contrib) in enumerate(members):
        mx, my = positions_row1[idx] if idx < 3 else positions_row2[idx - 3]

        # Light card background
        add_rect(sl, px(mx - 5), px(my - 5), Inches(4), Inches(1.7), LIGHT_GRAY)

        lines = [
            (role, 'Aptos', Pt(13), NAVY, True, False, PP_ALIGN.LEFT, Pt(8)),
            (f'Name \u2013 {name}', 'Aptos', Pt(12), NAVY, False, False, PP_ALIGN.LEFT, Pt(4)),
            (f'University \u2013 {uni}', 'Aptos', Pt(12), NAVY, False, False, PP_ALIGN.LEFT, Pt(4)),
            (f'DOB \u2013 {dob}', 'Aptos', Pt(12), NAVY, False, False, PP_ALIGN.LEFT, Pt(4)),
            (f'Contribution \u2013 {contrib}', 'Aptos', Pt(11), NAVY, False, False, PP_ALIGN.LEFT, Pt(4)),
        ]
        add_multi_para(sl, px(mx), px(my), Inches(3.8), Inches(1.6), lines)

    add_csic_logo(sl)
    add_slide_number(sl, 11)

    # ==============================================================
    # SLIDE 12: REFERENCES AND CITATIONS
    # ==============================================================
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, WHITE)
    add_slide_title(sl, 'References and Citations')

    refs = [
        '[1] NPCI, "UPI Product Statistics," 2024. npci.org.in/what-we-do/upi/product-statistics',
        '[2] RBI, "Master Direction on Digital Payment Security Controls," RBI/2020-21/74, 2021.',
        '[3] NPCI, "UPI Fraud Monitoring and Risk Management Guidelines," 2023.',
        '[4] S. Panigrahi et al., "Rule-Based and ML Methods for Fraud Detection," J. King Saud Univ., 2022.',
        '[5] E. Lopez-Rojas et al., "Applying AI and ML in Financial Services," IEEE Access, 2022.',
        '[6] G. Jambhrunkar et al., "MuleTrack: Temporal Learning for Money Mule Detection," IWANN, 2025.',
        '[7] D. Cheng et al., "GNNs for Financial Fraud Detection: A Review," arXiv:2411.05815, 2024.',
        '[8] M. Caglayan & S. Bahtiyar, "Money Laundering Detection with Node2Vec," Gazi Univ. J., 2022.',
        '[9] Z. Huang, "Money Mules Detection on Transaction Graphs," ACM GAIB, 2025.',
        '[10] Neo4j Inc., "Accelerate Fraud Detection with Graph Databases," Whitepaper, 2023.',
        '[11] F. T. Liu et al., "Isolation Forest," 8th IEEE Int. Conf. Data Mining, pp. 413-422, 2008.',
    ]

    ref_lines = []
    for ref in refs:
        ref_lines.append((ref, 'Aptos', Pt(11), NAVY, False, False, PP_ALIGN.LEFT, Pt(6)))

    add_multi_para(sl, px(25), px(78), Inches(12), Inches(3.8), ref_lines)

    # Project links
    add_text_box(sl, px(25), px(368), Inches(12), px(24),
                 'Project Links',
                 font_name='Aptos Display', font_size=Pt(16),
                 color=NAVY, bold=True)

    links = [
        ['Resource', 'Link'],
        ['GitHub Repository', '[Insert GitHub Repo URL]'],
        ['Live Deployed URL', '[Insert Deployed URL]'],
        ['Demo Video (YouTube)', '[Insert YouTube Link]'],
        ['API Documentation', '[Insert /docs URL when deployed]'],
    ]
    add_table(sl, px(25), px(395), Inches(8), links, font_size=Pt(11))

    add_csic_logo(sl)
    add_slide_number(sl, 12)

    # -- SAVE --
    prs.save(OUT_PATH)
    fsize = os.path.getsize(OUT_PATH) / 1024
    print(f"\n[+] Saved: {OUT_PATH}")
    print(f"    Size: {fsize:.1f} KB | Slides: {len(prs.slides)}")


if __name__ == '__main__':
    print("Generating FINAL_PRESENTATION.pptx (v2 - reviewed & perfected)...")
    build()
    print("Done!")
